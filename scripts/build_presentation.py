from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = PROJECT_ROOT / "FinSent_KJSSE_Template_Presentation.pptx"
OUTPUT_PATH = PROJECT_ROOT / "FinSent_KJSSE_Presentation_Pro.pptx"
ASSET_DIR = PROJECT_ROOT / "generated_ppt_assets"
SUMMARY_SHOWCASE_NAME = "summary_showcase.jpg"
NEWS_SHOWCASE_NAME = "news_showcase.jpg"



NAVY = RGBColor(13, 22, 39)
BLUE = RGBColor(32, 84, 145)
LIGHT_BLUE = RGBColor(224, 236, 248)
TEXT_DARK = RGBColor(31, 41, 55)
TEXT_LIGHT = RGBColor(243, 246, 251)
GREEN = RGBColor(51, 124, 86)
ORANGE = RGBColor(214, 132, 48)
RED = RGBColor(164, 53, 53)
GRAY = RGBColor(107, 114, 128)
SOFT = RGBColor(239, 243, 250)


def delete_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)  # type: ignore[attr-defined]
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)  # type: ignore[attr-defined]


def layout(prs: Presentation, name: str):
    for candidate in prs.slide_layouts:
        if candidate.name == name:
            return candidate
    raise ValueError(name)


def ensure_assets(asset_dir: Path = ASSET_DIR) -> None:
    asset_dir.mkdir(exist_ok=True)


def preprocess_image(src: Path, dest_name: str, crop: tuple[int, int, int, int] | None = None, asset_dir: Path = ASSET_DIR) -> Path:
    ensure_assets(asset_dir)
    image = Image.open(src).convert("RGB")
    if crop is not None:
        image = image.crop(crop)
    dest = asset_dir / dest_name
    image.save(dest, quality=92)
    return dest


def resolve_showcase_image(
    screenshot_path: Path | None,
    dest_name: str,
    crop: tuple[int, int, int, int],
    asset_dir: Path = ASSET_DIR,
) -> Path:
    if screenshot_path is not None:
        if not screenshot_path.exists():
            raise FileNotFoundError(f"Screenshot not found: {screenshot_path}")
        return preprocess_image(screenshot_path, dest_name, crop=crop, asset_dir=asset_dir)

    existing = asset_dir / dest_name
    if existing.exists():
        return existing
    raise FileNotFoundError(
        f"Missing {dest_name}. Provide a screenshot path or place the generated asset in {asset_dir}."
    )


def style_runs(paragraph, size: int = 20, color: RGBColor = TEXT_DARK, bold: bool = False) -> None:
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def clear_text_frame(tf) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP


def fill_textbox(tf, lines: list[tuple[str, int, int, RGBColor, bool]]) -> None:
    clear_text_frame(tf)
    first = True
    for text, level, size, color, bold in lines:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        paragraph.text = text
        paragraph.level = level
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(8 if level == 0 else 5)
        style_runs(paragraph, size=size, color=color, bold=bold)


def add_textbox(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    fill_textbox(box.text_frame, lines)
    return box


def add_card(slide, left, top, width, height, title: str, value: str, detail: str, fill=LIGHT_BLUE, value_color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BLUE
    tf = shape.text_frame
    fill_textbox(
        tf,
        [
            (title, 0, 14, GRAY, True),
            (value, 0, 26, value_color, True),
            (detail, 0, 12, TEXT_DARK, False),
        ],
    )
    return shape


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = BLUE
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def set_slide_title(slide, title: str) -> None:
    slide.shapes.title.text = title


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE"))
    slide.shapes.title.text = "FinSent"
    slide.placeholders[1].text = (
        "Provider-Driven Stock Intelligence Dashboard\n"
        "Saad Waghoo | K. J. Somaiya School of Engineering\n"
        "Department of Computer Engineering"
    )


def agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "OBJECT"))
    set_slide_title(slide, "Agenda")
    fill_textbox(
        slide.placeholders[1].text_frame,
        [
            ("1. Problem, motivation, and product thesis", 0, 24, TEXT_DARK, True),
            ("2. Exchange-first architecture and provider strategy", 0, 24, TEXT_DARK, True),
            ("3. Intelligence pipeline, scoring, and data honesty", 0, 24, TEXT_DARK, True),
            ("4. Dashboard walkthrough with real implementation highlights", 0, 24, TEXT_DARK, True),
            ("5. Engineering progress, limitations, and roadmap", 0, 24, TEXT_DARK, True),
        ],
    )


def executive_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "Executive Summary")
    add_textbox(
        slide,
        Inches(0.85),
        Inches(1.35),
        Inches(11.2),
        Inches(0.9),
        [
            (
                "FinSent is a multi-market stock intelligence platform that combines provider-backed quotes, recent ticker-specific news, structured LLM interpretation, and deterministic short-term scoring inside a polished Dash dashboard.",
                0,
                20,
                TEXT_DARK,
                False,
            )
        ],
    )
    add_card(slide, Inches(0.95), Inches(2.45), Inches(2.65), Inches(1.45), "Markets", "3 Exchanges", "US, NSE India, and BSE India")
    add_card(slide, Inches(3.9), Inches(2.45), Inches(2.65), Inches(1.45), "Coverage", "20 Symbols", "Seeded across tech, finance, energy, and consumer")
    add_card(slide, Inches(6.85), Inches(2.45), Inches(2.65), Inches(1.45), "Experience", "4 Pages", "Summary, News Impact, Compare, and Alerts")
    add_card(slide, Inches(9.8), Inches(2.45), Inches(2.0), Inches(1.45), "Progress", "~75%", "Implementation maturity vs planned scope", value_color=GREEN)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(4.35),
        Inches(10.9),
        Inches(1.55),
        [
            ("What makes it strong:", 0, 22, BLUE, True),
            ("Provider abstraction instead of one-off scripts", 0, 18, TEXT_DARK, False),
            ("Structured Gemini analysis with parse-safe fallback behavior", 0, 18, TEXT_DARK, False),
            ("Runtime badges that show live, stale, inferred, or unavailable states honestly", 0, 18, TEXT_DARK, False),
        ],
    )


def bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(layout(prs, "OBJECT"))
    set_slide_title(slide, title)
    fill_textbox(slide.placeholders[1].text_frame, [(bullet, 0, 22, TEXT_DARK, False) for bullet in bullets])


def two_col_slide(prs: Presentation, title: str, left_head: str, left_bullets: list[str], right_head: str, right_bullets: list[str]) -> None:
    slide = prs.slides.add_slide(layout(prs, "TWO_OBJECTS"))
    set_slide_title(slide, title)
    fill_textbox(
        slide.placeholders[1].text_frame,
        [(left_head, 0, 22, BLUE, True)] + [(item, 1, 18, TEXT_DARK, False) for item in left_bullets],
    )
    fill_textbox(
        slide.placeholders[2].text_frame,
        [(right_head, 0, 22, BLUE, True)] + [(item, 1, 18, TEXT_DARK, False) for item in right_bullets],
    )


def separator_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(layout(prs, "SECTION_HEADER"))
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def metric_grid_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "Coverage Snapshot")
    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.25),
        Inches(11.0),
        Inches(0.5),
        [("The current build is exchange-first, with a fixed registry that maps UI tickers to provider symbols cleanly.", 0, 18, TEXT_DARK, False)],
    )
    add_card(slide, Inches(0.95), Inches(1.95), Inches(3.55), Inches(1.4), "US Coverage", "8 symbols", "AAPL, AMZN, MSFT, NVDA, META, GOOGL, TSLA, JPM")
    add_card(slide, Inches(4.75), Inches(1.95), Inches(3.2), Inches(1.4), "NSE Coverage", "7 symbols", "RELIANCE, TCS, INFY, HDFCBANK, ICICIBANK, SBIN, ITC", fill=SOFT)
    add_card(slide, Inches(8.15), Inches(1.95), Inches(3.55), Inches(1.4), "BSE Coverage", "5 symbols", "RELIANCE, TCS, INFY, HDFCBANK, SBIN")
    add_card(slide, Inches(0.95), Inches(3.75), Inches(2.65), Inches(1.35), "Quote modes", "4 states", "live, stale, inferred, unavailable")
    add_card(slide, Inches(3.9), Inches(3.75), Inches(2.65), Inches(1.35), "News tiers", "2 levels", "provider-grade vs fallback-quality", fill=SOFT)
    add_card(slide, Inches(6.85), Inches(3.75), Inches(2.65), Inches(1.35), "Core entities", "4 stores", "quotes, bars, news, signal snapshots")
    add_card(slide, Inches(9.8), Inches(3.75), Inches(1.95), Inches(1.35), "Pages", "4", "summary, impact, compare, alerts", fill=SOFT)


def table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], col_widths: list[float]) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, title)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.55), Inches(1.45), Inches(11.8), Inches(4.9)).table
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            style_runs(paragraph, size=14, color=TEXT_LIGHT, bold=True)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 else SOFT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                style_runs(paragraph, size=13, color=TEXT_DARK, bold=False)


def runtime_modes_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "Runtime Modes and Data Honesty")
    add_card(slide, Inches(0.95), Inches(1.7), Inches(2.6), Inches(1.55), "Mode 1", "Market + News", "Fresh quotes and recent analyzed headlines are both active", fill=LIGHT_BLUE, value_color=GREEN)
    add_card(slide, Inches(3.75), Inches(1.7), Inches(2.6), Inches(1.55), "Mode 2", "Market-only fallback", "Quotes or bars are present, but headline support is weak or unavailable", fill=SOFT, value_color=ORANGE)
    add_card(slide, Inches(6.55), Inches(1.7), Inches(2.6), Inches(1.55), "Mode 3", "News-only fallback", "Headline coverage exists, but usable live market linkage is missing", fill=LIGHT_BLUE, value_color=ORANGE)
    add_card(slide, Inches(9.35), Inches(1.7), Inches(2.3), Inches(1.55), "Mode 4", "Unavailable", "The app refuses to fake live values when providers fail", fill=SOFT, value_color=RED)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(3.7),
        Inches(10.9),
        Inches(1.55),
        [
            ("This honesty layer is a major difference between FinSent and a typical student-demo dashboard.", 0, 20, BLUE, True),
            ("The UI surfaces price source, news provider, provider tier, freshness age, and quality badges so the user always knows what is grounded versus inferred.", 0, 18, TEXT_DARK, False),
        ],
    )


def architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "BLANK"))
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(5.2), Inches(0.55), [("High-Level Architecture", 0, 26, TEXT_DARK, True)])
    boxes = [
        add_card(slide, Inches(0.7), Inches(2.35), Inches(1.7), Inches(1.0), "User Layer", "User", "Ticker selection and page navigation", fill=SOFT),
        add_card(slide, Inches(2.65), Inches(2.35), Inches(2.0), Inches(1.0), "Frontend", "Dash UI", "Summary, News Impact, Compare, Alerts"),
        add_card(slide, Inches(4.95), Inches(2.35), Inches(1.9), Inches(1.0), "State", "View Model", "Loads and normalizes data for figures and cards", fill=SOFT),
        add_card(slide, Inches(7.15), Inches(1.55), Inches(2.25), Inches(1.75), "Providers", "Market + News", "Polygon, Marketaux, Kite plan, fallback web"),
        add_card(slide, Inches(7.15), Inches(3.55), Inches(2.25), Inches(1.75), "Intelligence", "Gemini + Heuristics", "Structured interpretation with deterministic signal logic", fill=SOFT),
        add_card(slide, Inches(9.75), Inches(2.35), Inches(2.0), Inches(1.0), "Persistence", "SQLite / SQLAlchemy", "Quotes, bars, news, and signal snapshots"),
    ]
    add_connector(slide, boxes[0].left + boxes[0].width, boxes[0].top + boxes[0].height // 2, boxes[1].left, boxes[1].top + boxes[1].height // 2)
    add_connector(slide, boxes[1].left + boxes[1].width, boxes[1].top + boxes[1].height // 2, boxes[2].left, boxes[2].top + boxes[2].height // 2)
    add_connector(slide, boxes[2].left + boxes[2].width, boxes[2].top + boxes[2].height // 2, boxes[3].left, boxes[3].top + boxes[3].height // 2)
    add_connector(slide, boxes[3].left + boxes[3].width // 2, boxes[3].top + boxes[3].height, boxes[4].left + boxes[4].width // 2, boxes[4].top)
    add_connector(slide, boxes[4].left + boxes[4].width, boxes[4].top + boxes[4].height // 2, boxes[5].left, boxes[5].top + boxes[5].height // 2)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.85),
        Inches(10.9),
        Inches(0.65),
        [("Key principle: providers fetch raw data, Gemini interprets the data, and the app owns the final composite signal.", 0, 18, BLUE, True)],
    )


def pipeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "BLANK"))
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(6), Inches(0.55), [("Runtime Pipeline", 0, 26, TEXT_DARK, True)])
    steps = [
        "1. User selects exchange and symbol",
        "2. Symbol registry resolves the provider symbol",
        "3. Market provider fetches quote snapshot and price bars",
        "4. News provider fetches recent headlines",
        "5. Gemini returns structured JSON or heuristic fallback is used",
        "6. Signal engine computes deterministic composite score",
        "7. View model renders cards, charts, tables, and alerts",
    ]
    previous = None
    top = Inches(1.35)
    for idx, step in enumerate(steps):
        box = add_card(slide, Inches(1.2), top + Inches(idx * 0.68), Inches(9.8), Inches(0.5), "", step, "", fill=LIGHT_BLUE if idx % 2 == 0 else SOFT, value_color=TEXT_DARK)
        if previous is not None:
            add_connector(slide, previous.left + previous.width // 2, previous.top + previous.height, box.left + box.width // 2, box.top)
        previous = box


def screenshot_showcase_slide(prs: Presentation, title: str, image_path: Path, callouts: list[tuple[str, str]]) -> None:
    slide = prs.slides.add_slide(layout(prs, "BLANK"))
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(5.5), Inches(0.55), [(title, 0, 26, TEXT_DARK, True)])
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.15), Inches(7.6), Inches(5.25))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(248, 250, 253)
    frame.line.color.rgb = BLUE
    slide.shapes.add_picture(str(image_path), Inches(0.92), Inches(1.3), width=Inches(7.25))
    top = Inches(1.35)
    for idx, (head, body) in enumerate(callouts):
        add_card(slide, Inches(8.7), top + Inches(idx * 1.52), Inches(3.05), Inches(1.2), head, "", body, fill=LIGHT_BLUE if idx % 2 == 0 else SOFT, value_color=BLUE)


def compare_alerts_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "Compare and Alerts Layers")
    add_card(slide, Inches(0.95), Inches(1.55), Inches(5.2), Inches(1.55), "Compare page", "Peer benchmarking", "Relative performance, signal leadership, news intensity, and confidence comparison across selected symbols")
    add_card(slide, Inches(6.45), Inches(1.55), Inches(5.2), Inches(1.55), "Alerts page", "Monitoring layer", "Recent signal shifts, active alerts, sector mood, and a plain-language Buy / Watch / Avoid read")
    add_textbox(
        slide,
        Inches(0.95),
        Inches(3.65),
        Inches(10.75),
        Inches(1.65),
        [
            ("These pages make the system feel less like a static analysis sheet and more like an ongoing stock-intelligence workspace.", 0, 20, BLUE, True),
            ("Even where live India prices are blocked, the UX still explains what is missing instead of quietly degrading into fake values.", 0, 18, TEXT_DARK, False),
        ],
    )


def trust_layer_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "UX Trust Layer")
    items = [
        ("Price source", "Shows whether the price came from Polygon, Kite, or fallback logic"),
        ("News provider", "Surfaces Polygon, Marketaux, or fallback web source directly in the UI"),
        ("Provider tier", "Marks provider-grade versus fallback-quality coverage"),
        ("Freshness age", "Makes staleness visible instead of hiding it"),
        ("Parse status", "Explains when Gemini succeeded, failed, or fell back to heuristics"),
        ("Mode", "Communicates Market + News, Market-only fallback, News-only fallback, or Unavailable"),
    ]
    left = Inches(0.95)
    top = Inches(1.55)
    for idx, (head, body) in enumerate(items):
        x = left + Inches((idx % 2) * 5.45)
        y = top + Inches((idx // 2) * 1.35)
        add_card(slide, x, y, Inches(5.0), Inches(1.05), head, "", body, fill=LIGHT_BLUE if idx % 2 == 0 else SOFT, value_color=BLUE)


def built_modules_slide(prs: Presentation) -> None:
    rows = [
        ["market_providers.py", "Quote snapshots, price bars, provider routing", "Built"],
        ["news_providers.py", "Polygon, Marketaux, and fallback news ingestion", "Built"],
        ["llm_analyzers.py", "Gemini structured analysis plus heuristic fallback", "Built"],
        ["signal_engine.py", "Deterministic composite scoring", "Built"],
        ["intelligence_service.py", "Orchestrates provider outputs into stored insights", "Built"],
        ["dashboard/view_model.py", "Transforms persistence into cards, charts, and status labels", "Built"],
    ]
    table_slide(prs, "What Was Actually Built", ["Module", "Responsibility", "Status"], rows, [3.8, 6.3, 1.5])


def fixes_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "OBJECT"))
    set_slide_title(slide, "Major Engineering Fixes Already Solved")
    bullets = [
        "Gemini quota and malformed JSON no longer break the UX; heuristic fallback fills confidence and impact safely.",
        "News Impact tables now show all recent headlines, not just the small subset with direct price overlap.",
        "Failed quote rows no longer override usable older rows during auto-refresh.",
        "Ancient historical closes are blocked from appearing as fake current prices.",
        "Charts can fall back to estimated impact when direct event-linkage is sparse.",
    ]
    fill_textbox(slide.placeholders[1].text_frame, [(bullet, 0, 20, TEXT_DARK, False) for bullet in bullets])


def limitations_slide(prs: Presentation) -> None:
    rows = [
        ["India live quote source", "NSE/BSE prices remain unavailable without Kite or a replacement provider", "Complete Kite setup or replace India quote path"],
        ["Polygon plan limits", "US snapshot quality depends on which endpoints are allowed by the API plan", "Use last-trade / previous-close fallback and label it honestly"],
        ["Gemini quota", "Heavy repeated refreshes can force heuristic analysis", "Cache article hashes and limit analysis budget per refresh"],
        ["Demo stability", "Some provider outages can leave the UI sparse", "Keep screenshots and runtime-state explanations ready"],
    ]
    table_slide(prs, "Current Limitations and Mitigations", ["Constraint", "Effect on product", "Mitigation"], rows, [2.6, 4.9, 4.0])


def progress_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "Progress vs Planned Timeline")
    add_textbox(slide, Inches(0.95), Inches(1.35), Inches(11.0), Inches(0.55), [("By schedule the project is in Phase 5, but implementation maturity is already closer to late Phase 6 / early Phase 7.", 0, 19, TEXT_DARK, False)])
    bars = [
        ("Scheduled timeline completion", 45, BLUE),
        ("Actual implementation maturity", 75, GREEN),
    ]
    for idx, (label, pct, color) in enumerate(bars):
        y = Inches(2.2 + idx * 1.45)
        add_textbox(slide, Inches(0.95), y, Inches(4.0), Inches(0.35), [(f"{label}: {pct}%", 0, 20, TEXT_DARK, True)])
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), y + Inches(0.45), Inches(10.6), Inches(0.3))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(231, 236, 243)
        bg.line.color.rgb = RGBColor(231, 236, 243)
        fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), y + Inches(0.45), Inches(10.6 * pct / 100), Inches(0.3))
        fg.fill.solid()
        fg.fill.fore_color.rgb = color
        fg.line.color.rgb = color
    add_textbox(
        slide,
        Inches(0.95),
        Inches(5.0),
        Inches(10.7),
        Inches(1.0),
        [
            ("The calendar says ~45% complete; the actual build is around 75% complete because core architecture, multi-page UI, persistence, provider routing, and reliability fixes are already in place.", 0, 18, BLUE, True)
        ],
    )


def portfolio_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE_ONLY"))
    set_slide_title(slide, "Why This Feels Portfolio-Grade")
    add_card(slide, Inches(0.95), Inches(1.6), Inches(2.55), Inches(1.45), "Architecture", "Modular", "Providers, LLM analysis, signal engine, persistence, and UI are cleanly separated")
    add_card(slide, Inches(3.75), Inches(1.6), Inches(2.55), Inches(1.45), "LLM role", "Constrained", "Gemini interprets fetched news instead of inventing unsupported market facts", fill=SOFT)
    add_card(slide, Inches(6.55), Inches(1.6), Inches(2.55), Inches(1.45), "Scoring", "Deterministic", "Final bullish / bearish signal is computed in application logic", fill=LIGHT_BLUE)
    add_card(slide, Inches(9.35), Inches(1.6), Inches(2.3), Inches(1.45), "UX", "Honest", "Live, stale, inferred, and unavailable states are visible", fill=SOFT)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(3.65),
        Inches(10.9),
        Inches(1.8),
        [
            ("This is no longer a static CSV demo.", 0, 22, BLUE, True),
            ("It behaves like a production-minded intelligence product: provider driven, failure aware, and explicit about uncertainty.", 0, 18, TEXT_DARK, False),
            ("That separation between data ingestion, AI interpretation, and deterministic scoring is the strongest part of the build.", 0, 18, TEXT_DARK, False),
        ],
    )


def roadmap_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TWO_OBJECTS"))
    set_slide_title(slide, "Next Steps")
    fill_textbox(
        slide.placeholders[1].text_frame,
        [
            ("Near-term", 0, 22, BLUE, True),
            ("Finish India live quote path with Kite or a replacement provider", 1, 18, TEXT_DARK, False),
            ("Capture stable dashboard screenshots for the final viva", 1, 18, TEXT_DARK, False),
            ("Run final end-to-end refresh checks across AAPL, NVDA, and India symbols", 1, 18, TEXT_DARK, False),
        ],
    )
    fill_textbox(
        slide.placeholders[2].text_frame,
        [
            ("Longer-term", 0, 22, BLUE, True),
            ("Move persistence from SQLite to PostgreSQL", 1, 18, TEXT_DARK, False),
            ("Add authentication and background refresh scheduling", 1, 18, TEXT_DARK, False),
            ("Turn the dashboard into a cloud-deployable monitoring product", 1, 18, TEXT_DARK, False),
        ],
    )


def closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(layout(prs, "TITLE"))
    slide.shapes.title.text = "FinSent"
    slide.placeholders[1].text = "Provider-driven, explainable, and production-minded.\nThank you."


def build(
    template_path: Path = TEMPLATE_PATH,
    output_path: Path = OUTPUT_PATH,
    asset_dir: Path = ASSET_DIR,
    summary_screenshot: Path | None = None,
    news_screenshot: Path | None = None,
) -> Path:
    ensure_assets(asset_dir)
    summary_image = resolve_showcase_image(summary_screenshot, SUMMARY_SHOWCASE_NAME, crop=(150, 40, 3170, 1900), asset_dir=asset_dir)
    news_image = resolve_showcase_image(news_screenshot, NEWS_SHOWCASE_NAME, crop=(120, 120, 3190, 1750), asset_dir=asset_dir)

    prs = Presentation(str(template_path))
    delete_all_slides(prs)

    title_slide(prs)
    agenda_slide(prs)
    executive_summary_slide(prs)
    bullet_slide(
        prs,
        "Problem Statement",
        [
            "Retail traders typically check prices, news, analyst updates, and social commentary across disconnected tools.",
            "Most dashboards either show static historical values or use AI summaries without clear source ownership.",
            "The real challenge is not just sentiment classification; it is combining market context, data freshness, and explainability in one system.",
        ],
    )
    two_col_slide(
        prs,
        "Where Generic Dashboards Fail",
        "Typical student-demo weaknesses",
        [
            "Hardcoded tickers and CSV-driven values",
            "No distinction between live data and stale data",
            "LLM treated like a data source instead of an analyst",
        ],
        "What FinSent does differently",
        [
            "Exchange-first flow with provider routing",
            "Runtime honesty through source, freshness, and quality badges",
            "Deterministic signal engine on top of structured AI output",
        ],
    )
    bullet_slide(
        prs,
        "Product Thesis",
        [
            "FinSent should feel like a real stock-intelligence workspace, not a notebook demo.",
            "The product should tell the user what the market is doing, what the news is saying, how strong that evidence is, and whether the system trusts its own conclusion.",
            "That is why provider abstraction, data honesty, and explanation quality matter as much as raw sentiment accuracy.",
        ],
    )

    separator_slide(prs, "Platform Scope", "Coverage, providers, and runtime trust model")
    metric_grid_slide(prs)
    table_slide(
        prs,
        "Provider Matrix",
        ["Layer", "US path", "India path", "What the app does"],
        [
            ["Quote snapshot", "Polygon", "Kite (planned / credential-gated)", "Fetches price, timestamp, freshness, and liquidity context"],
            ["Historical bars", "Polygon aggregates", "Kite historical bars", "Drives timelines, overlap analysis, and price windows"],
            ["News ingestion", "Polygon News", "Marketaux", "Normalizes article metadata and relevance"],
            ["LLM analysis", "Gemini", "Gemini", "Returns structured relevance, tone, confidence, impact, catalyst"],
            ["Fallback", "last-trade / prev-close / web news", "web news / unavailable", "Protects the UX without pretending stale data is live"],
        ],
        [2.1, 2.5, 2.8, 4.2],
    )
    runtime_modes_slide(prs)

    separator_slide(prs, "Architecture", "How FinSent is assembled as a real system")
    architecture_slide(prs)
    pipeline_slide(prs)
    table_slide(
        prs,
        "Persistence and Data Model",
        ["Store", "Purpose", "Key fields"],
        [
            ["Quote snapshots", "Latest market state", "ticker, exchange, current price, bid/ask, provider, freshness, quality"],
            ["Price bars", "Historical movement", "open, high, low, close, volume, interval, market timestamp"],
            ["Normalized news", "Headline evidence", "title, source, url, provider, sentiment, confidence, parse status"],
            ["Signal snapshots", "Decision layer", "composite score, label, action bias, mode, final reason"],
        ],
        [2.1, 2.8, 6.3],
    )
    two_col_slide(
        prs,
        "Intelligence Pipeline",
        "Structured article analysis",
        [
            "relevant / not relevant",
            "bullish / bearish / neutral",
            "confidence and impact strength",
            "time horizon and catalyst tag",
        ],
        "Reliability controls",
        [
            "JSON parsing validation",
            "Heuristic fallback for quota or parse failures",
            "Cached article hashes to avoid repeat analysis",
            "Quality badge downgraded when AI output is inferred",
        ],
    )
    two_col_slide(
        prs,
        "Composite Signal Engine",
        "Inputs blended by app logic",
        [
            "aggregate news sentiment",
            "article recency and relevance",
            "market pressure and price movement",
            "volume and liquidity context",
        ],
        "Outputs visible in the UI",
        [
            "composite score",
            "bullish / neutral / bearish label",
            "signal confidence",
            "plain-language explanation",
        ],
    )

    separator_slide(prs, "Product Experience", "What the dashboard actually shows")
    screenshot_showcase_slide(
        prs,
        "Summary Page Walkthrough",
        summary_image,
        [
            ("Top workspace bar", "Exchange and ticker selection route the request to the correct provider stack."),
            ("Signal cards", "Current price, composite score, confidence, and buy/sell proxy are shown together."),
            ("Explanation box", "The system explains why the current signal looks the way it does in short readable language."),
        ],
    )
    screenshot_showcase_slide(
        prs,
        "News Impact Page Walkthrough",
        news_image,
        [
            ("Impact map", "Plots sentiment against estimated impact when direct price overlap is available."),
            ("Current window", "Compresses headline count, confidence, and impact summary."),
            ("Headline table", "Shows provider, source, headline, sentiment, confidence, and explanation row by row."),
        ],
    )
    compare_alerts_slide(prs)
    trust_layer_slide(prs)

    separator_slide(prs, "Engineering Quality", "Why the implementation is stronger than a static prototype")
    built_modules_slide(prs)
    fixes_slide(prs)
    limitations_slide(prs)
    progress_slide(prs)

    separator_slide(prs, "Value and Positioning", "Why this is already stronger than a typical academic build")
    portfolio_slide(prs)
    roadmap_slide(prs)
    closing_slide(prs)

    prs.save(str(output_path))
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the FinSent presentation from project-local assets.")
    parser.add_argument("--template", type=Path, default=TEMPLATE_PATH, help="PowerPoint template path.")
    parser.add_argument("--output", type=Path, default=OUTPUT_PATH, help="Output presentation path.")
    parser.add_argument("--asset-dir", type=Path, default=ASSET_DIR, help="Directory for generated presentation images.")
    parser.add_argument("--summary-screenshot", type=Path, default=None, help="Optional raw summary-page screenshot to crop.")
    parser.add_argument("--news-screenshot", type=Path, default=None, help="Optional raw news-impact screenshot to crop.")
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()
    path = build(
        template_path=args.template,
        output_path=args.output,
        asset_dir=args.asset_dir,
        summary_screenshot=args.summary_screenshot,
        news_screenshot=args.news_screenshot,
    )
    print(path)
